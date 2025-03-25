from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from config.settings import OPENAI_API_KEY, BLOB_CONNECTION_STR, CONTAINER_NAME
from utils.loader import parse_response
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import base64
from botbuilder.schema import Attachment, ActivityTypes, Activity
from azure.storage.blob import BlobServiceClient, ContentSettings

proposal_prompt = PromptTemplate(
    input_variables=["user_message"],
    template="""
You are a proposal presentation generator.
Draft a PowerPoint presentation outline for a proposal based on the user input.
The presentation should include:
- A Title Slide (with a proposal title and an optional subtitle)
- An Agenda Slide (listing main topics)
- Two slides with mocked relevant use cases
- A Conclusion Slide

User Input: {user_message}

Output the presentation outline in the following format:

Slide 1: Title Slide
- Title: [Your Proposal Title]
- Subtitle: [Optional Subtitle]

Slide 2: Agenda
- [Agenda bullet point 1]
- [Agenda bullet point 2]
- [Agenda bullet point 3]

Slide 3: Use Case 1
- [Description of Use Case 1]

Slide 4: Use Case 2
- [Description of Use Case 2]

Slide 5: Conclusion
- [Key takeaway or next steps]

Please provide the full draft outline.
"""
)

def generate_ppt_from_outline(outline: str) -> bytes:
    prs = Presentation()
    slides_text = outline.strip().split("\n\n")
    for slide_text in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = textbox.text_frame
        tf.text = ""
        for line in slide_text.split("\n"):
            p = tf.add_paragraph()
            p.text = line.strip()
    output = BytesIO()
    prs.save(output)
    ppt_data = output.getvalue()
    output.close()
    return ppt_data

class ProposalHandler:
    def __init__(self, llm):
        self.chain = proposal_prompt | llm

    def upload_file_to_blob(self, blob_name: str, file_bytes: bytes) -> str:
        blob_service_client = BlobServiceClient.from_connection_string(BLOB_CONNECTION_STR)
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)

        container_client.upload_blob(
            name=blob_name,
            data=file_bytes,
            overwrite=True,
            content_settings=ContentSettings(
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        )

        account_url = blob_service_client.primary_endpoint
        return f"{account_url}/{CONTAINER_NAME}/{blob_name}"

    async def handle(self, user_input: str) -> dict:
        proposal_text = await self.chain.ainvoke({"user_message":user_input})
        proposal_text = parse_response(proposal_text)
        ppt_data = generate_ppt_from_outline(proposal_text)
        base64_data = base64.b64encode(ppt_data).decode("utf-8")

        blob_name = "proposal.pptx"
        content_url = self.upload_file_to_blob(blob_name, ppt_data)

        attachment = Attachment(
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            content=base64_data,
            content_url=content_url,
            name=blob_name
        )

        activity = Activity(
            type=ActivityTypes.message,
            text="Here is your proposal presentation draft attached.",
            attachments=[attachment]
        )

        return activity