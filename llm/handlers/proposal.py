from langchain_openai import ChatOpenAI
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from config.settings import OPENAI_API_KEY, BLOB_CONNECTION_STR, CONTAINER_NAME
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import base64
from botbuilder.schema import Attachment, ActivityTypes, Activity
from azure.storage.blob import BlobServiceClient, ContentSettings

llm = ChatOpenAI(model_name="gpt-4", openai_api_key=OPENAI_API_KEY, temperature=0.7)

proposal_prompt = PromptTemplate(
    input_variables=["user_message"],
    template="""
You are a proposal presentation generator.
Draft a PowerPoint presentation outline for a proposal based on the user input.
The presentation should include:
- A Title Slide (with a proposal title and an optional subtitle)
- An Agenda Slide (listing main topics)
- Two slides with mocked relevant use cases
- A Conclusion Slide

User Input: {user_message}

Output the presentation outline in the following format:

Slide 1: Title Slide
- Title: [Your Proposal Title]
- Subtitle: [Optional Subtitle]

Slide 2: Agenda
- [Agenda bullet point 1]
- [Agenda bullet point 2]
- [Agenda bullet point 3]

Slide 3: Use Case 1
- [Description of Use Case 1]

Slide 4: Use Case 2
- [Description of Use Case 2]

Slide 5: Conclusion
- [Key takeaway or next steps]

Please provide the full draft outline.
"""
)

def generate_ppt_from_outline(outline: str) -> bytes:
    prs = Presentation()
    # Split the outline into separate slide blocks by double newlines.
    slides_text = outline.strip().split("\n\n")
    for slide_text in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank layout.
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tf = textbox.text_frame
        # Clear the default paragraph that is automatically added.
        tf.text = ""
        for line in slide_text.split("\n"):
            p = tf.add_paragraph()
            p.text = line.strip()
    output = BytesIO()
    prs.save(output)
    ppt_data = output.getvalue()
    output.close()
    return ppt_data

class ProposalHandler:
    def __init__(self):
        self.chain = LLMChain(llm=llm, prompt=proposal_prompt)

    def upload_file_to_blob(blob_name: str) -> str:
        blob_service_client = BlobServiceClient.from_connection_string(BLOB_CONNECTION_STR)
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
        blob_name = "proposal.pptx"
        with open(blob_name, "rb") as data:
            container_client.upload_blob(
                name=blob_name,
                data=data,
                overwrite=True,
                content_settings=ContentSettings(
                    content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            )

        account_url = blob_service_client.primary_endpoint
        file_url = f"{account_url}/{CONTAINER_NAME}/{blob_name}"
        return file_url

    async def handle(self, user_input: str) -> dict:
        # Generate the proposal outline using the LLM.
        proposal_text = await self.chain.arun(user_message=user_input)
        # Generate a PowerPoint file from the outline.
        ppt_data = generate_ppt_from_outline(proposal_text)
        base64_data = base64.b64encode(ppt_data).decode('utf-8')
        name = "proposal.pptx"
        content_url = self.upload_file_to_blob(name)

        attachment = Attachment(
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            content=base64_data,
            content_url=content_url,
            name="proposal.pptx"
        )

        activity = Activity(
            type=ActivityTypes.message,
            text="Here is your proposal presentation draft attached.",
            attachments=[attachment]
        )

        return activity
